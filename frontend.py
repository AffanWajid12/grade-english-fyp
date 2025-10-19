import streamlit as st
import requests
import PyPDF2
import io
from pptx import Presentation
import json
from datetime import datetime
# --- Add this after imports (once) ---
st.markdown(
    """
    <style>
    /* Box styles for nested questions */
    .aiq-box {
      border: 1px solid #dfe6ee;
      border-radius: 10px;
      padding: 10px;
      margin-bottom: 12px;
      background: #ffffff;
      box-shadow: 0 1px 2px rgba(16,24,40,0.03);
    }
    .aiq-box.depth-1 { margin-left: 0px; }
    .aiq-box.depth-2 { margin-left: 18px; background: #fbfdff; }
    .aiq-box.depth-3 { margin-left: 34px; background: #fbfbfb; }
    .aiq-box.depth-4 { margin-left: 52px; background: #fafafa; }
    .aiq-box .aiq-header { font-weight: 700; margin-bottom: 8px; color: #0f172a; }
    .aiq-small { font-size: 13px; color: #475569; }
    </style>
    """,
    unsafe_allow_html=True
)

# ================== CONFIGURATION ==================
STRICTNESS_LEVELS = {
    1: "Be extremely lenient. Award points for any attempt.",
    2: "Be very lenient. Award partial credit generously.",
    3: "Be lenient. The answer can be partially incorrect.",
    4: "Be slightly lenient. Minor inaccuracies are acceptable.",
    5: "Be balanced. A fair mix of strictness and leniency.",
    6: "Be slightly strict. The answer must be mostly correct.",
    7: "Be strict. Deduct points for any inaccuracies.",
    8: "Be very strict. The answer must be precise and fully correct.",
    9: "Be extremely strict. The answer must be perfect.",
    10: "Be unforgiving. Any deviation results in zero points."
}

BACKEND_URL_GRADE = "http://127.0.0.1:5000/grade"
BACKEND_URL_RUBRIC_SINGLE = "http://127.0.0.1:5000/generate_question_rubric"
BACKEND_URL_RUBRIC_BULK = "http://127.0.0.1:5000/generate_rubrics_bulk"
BACKEND_URL_DELETE_RUBRIC = "http://127.0.0.1:5000/delete_rubric"

# ================== HELPERS ==================
def extract_text_from_files(uploaded_files):
    """Extract readable text from PDF, TXT, and PPTX files."""
    full_text = []
    if not uploaded_files:
        return ""
    for uploaded_file in uploaded_files:
        try:
            if uploaded_file.type == "text/plain":
                full_text.append(uploaded_file.getvalue().decode("utf-8"))
            elif uploaded_file.type == "application/pdf":
                pdf_reader = PyPDF2.PdfReader(uploaded_file)
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text:
                        full_text.append(text)
            elif "presentationml" in uploaded_file.type:
                prs = Presentation(io.BytesIO(uploaded_file.getvalue()))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            full_text.append(shape.text)
        except Exception as e:
            st.error(f"⚠️ Error processing file {uploaded_file.name}: {e}")
    return "\n".join(full_text)

def get_item_by_path(path):
    """
    Given a path list like [0,1,2], return the item reference inside st.session_state.exam_questions.
    This allows us to update the item in-place.
    """
    cur = st.session_state.exam_questions
    for idx in path:
        # if cur is a list (top-level), index directly
        if isinstance(cur, list):
            cur = cur[idx]
        else:
            # cur is a dict (question) so go into its 'parts'
            cur = cur.get("parts", [])[idx]
    return cur

# ================== PAGE SETUP ==================
st.set_page_config(layout="wide", page_title="AI Grading Assistant")
st.title("📝 AI Grading Assistant v12.2")

if "exam_questions" not in st.session_state:
    st.session_state.exam_questions = []

# ================== SIDEBAR ==================
with st.sidebar:
    st.header("⚙️ Configuration")

    use_strictness = st.toggle("Enable Strictness Level", value=True)
    strictness_instruction = ""
    if use_strictness:
        level = st.slider("Strictness Level", 1, 10, 6)
        strictness_instruction = STRICTNESS_LEVELS[level]
        st.info(f"**AI Instruction:** {strictness_instruction}")

    additional_instructions = st.text_area(
        "Custom AI Command (Optional)",
        height=150,
        help="Overrides strictness level if filled."
    )

    uploaded_files = st.file_uploader(
        "Upload Course Materials (Optional)",
        type=["pdf", "txt", "pptx"],
        accept_multiple_files=True
    )

# ================== QUESTION MANAGEMENT ==================
def delete_question_part(path):
    data = st.session_state.exam_questions
    for idx in path[:-1]:
        data = data[idx]["parts"]
    data.pop(path[-1])

def add_sub_part(path):
    data = st.session_state.exam_questions
    for idx in path:
        data = data[idx]["parts"]
    data.append({
        "question": "",
        "answer": "",
        "points": 1,
        "type": "short_question",
        "parts": [],
        "rubric": None,
        "rubric_cols": 3,
        "suggested_criteria": ""
    })

def generate_rubric_for_question(path, force=False):
    """
    Generate rubric for a specific question identified by path list.
    If force=True, we overwrite local rubric regardless of existing value.
    """
    try:
        item = get_item_by_path(path)
    except Exception as e:
        st.error(f"Internal error locating question by path {path}: {e}")
        return

    question_text = item.get("question", "")
    points = item.get("points", 0)
    cols = int(item.get("rubric_cols", 3))
    suggested_criteria = item.get("suggested_criteria", "").strip() or None

    if not question_text or not points:
        st.warning("Please provide question text and points before generating a rubric.")
        return

    if item.get("rubric") and not force:
        st.info("Rubric already exists for this question. Use 'Regenerate' to overwrite.")
        return

    payload = {
        "question": question_text,
        "points": points,
        "columns": cols,
        "criteria": [],  # let AI generate
        "path": path  # send path so backend returns path_key for reliable mapping
    }

    if suggested_criteria:
        payload["criteria"] = [c.strip() for c in suggested_criteria.split(",") if c.strip()]

    with st.spinner("Generating rubric using AI..."):
        try:
            resp = requests.post(BACKEND_URL_RUBRIC_SINGLE, json=payload, timeout=120)
            resp.raise_for_status()
            data = resp.json()
            # API returns { "path_key": "q-0-1", "rubric": [...] }
            rubric_matrix = data.get("rubric")
            if rubric_matrix is None:
                st.error("Backend returned no rubric.")
                return
            # Set rubric into our exam_questions item
            item["rubric"] = rubric_matrix
            st.session_state.exam_questions = st.session_state.exam_questions  # trigger session persistence
            st.success("Rubric generated and attached to the question ✅")
        except requests.exceptions.RequestException as e:
            st.error(f"Failed to generate rubric: {e}")
        except Exception as e:
            st.error(f"Unexpected error: {e}")

def regenerate_rubric_for_question(path):
    """Force regenerate a rubric for a single question."""
    # call generate_rubric_for_question with force True (we overwrite locally)
    generate_rubric_for_question(path, force=True)

def delete_rubric_for_question(path):
    """
    Call backend delete endpoint (for compatibility). Then clear rubric locally.
    Backend returns deleted path_key; we still clear locally because frontend holds canonical state.
    """
    # Build path_key
    path_key = "q-" + "-".join(map(str, path))
    payload = {"path_key": path_key}
    try:
        with st.spinner("Deleting rubric..."):
            resp = requests.post(BACKEND_URL_DELETE_RUBRIC, json=payload, timeout=30)
            resp.raise_for_status()
            data = resp.json()
            deleted = data.get("deleted")
            # Clear local rubric
            item = get_item_by_path(path)
            if item.get("rubric"):
                item["rubric"] = None
            st.session_state.exam_questions = st.session_state.exam_questions
            st.success(f"Rubric deleted for {path_key}.")
    except requests.exceptions.RequestException as e:
        st.error(f"Failed to delete rubric: {e}")
    except Exception as e:
        st.error(f"Unexpected error: {e}")

def generate_rubrics_for_all(force=False, columns=3, timeout_per_call=60, max_workers=6):
    """
    Generate rubrics for all leaf questions by calling backend /generate_rubrics_bulk.
    Updates st.session_state.exam_questions in-place with returned rubrics.
    """
    payload = {
        "student_exam": st.session_state.exam_questions,
        "columns": int(columns),
        "force": bool(force),
        "timeout_per_call": int(timeout_per_call),
        "max_workers": int(max_workers)
    }

    with st.spinner("Generating rubrics for all leaf questions. This may take a while..."):
        try:
            resp = requests.post(BACKEND_URL_RUBRIC_BULK, json=payload, timeout=300)
            resp.raise_for_status()
            data = resp.json()
            generated = data.get("generated", {})
            skipped = data.get("skipped", [])

            # Apply generated rubrics into exam_questions by path_key -> path list conversion
            applied = 0
            errors = 0
            for pk, rubric_or_err in generated.items():
                # pk like "q-0-1"
                if not pk.startswith("q-"):
                    continue
                parts = pk.split("-")[1:]
                try:
                    path = [int(p) for p in parts]
                except ValueError:
                    continue
                # if rubric_or_err is dict with "error", show as error
                if isinstance(rubric_or_err, dict) and "error" in rubric_or_err:
                    errors += 1
                    # optionally attach an error marker to question for UI
                    try:
                        item = get_item_by_path(path)
                        item["rubric_error"] = rubric_or_err["error"]
                    except Exception:
                        pass
                    continue

                try:
                    item = get_item_by_path(path)
                    item["rubric"] = rubric_or_err
                    applied += 1
                except Exception:
                    errors += 1

            st.session_state.exam_questions = st.session_state.exam_questions
            msg = f"Generated rubrics applied: {applied}. Skipped (existing): {len(skipped)}. Errors: {errors}."
            st.success(msg)
        except requests.exceptions.RequestException as e:
            st.error(f"Bulk generation failed: {e}")
        except Exception as e:
            st.error(f"Unexpected error during bulk generation: {e}")

# ================== DISPLAY QUESTIONS (stable box-in-box using expanders) ==================
def display_questions(questions, path_prefix=[], depth=1):
    """
    Render questions as nested boxes using st.expander and st.container.
    Each question shows only the editable widgets once; sub-parts are rendered inside
    the parent's expander so they appear visually nested.
    """
    for i, item in enumerate(questions):
        path = path_prefix + [i]
        key_prefix = "-".join(map(str, path))
        # Use an expander as the visual box (expanded by default)
        # Use a small label so the UI shows a clean border and body for widgets.
        with st.expander(label="", expanded=True):
            # Top row: question input + add / delete buttons
            cols_top = st.columns([8, 1, 1])
            item["question"] = cols_top[0].text_input("Question", value=item.get("question", ""), key=f"q_{key_prefix}")
            cols_top[1].button("➕ Sub-part", key=f"add_{key_prefix}", on_click=add_sub_part, args=(path,))
            cols_top[2].button("❌ Delete", key=f"del_{key_prefix}", on_click=delete_question_part, args=(path,))

            # Meta: type and points
            cols_meta = st.columns([3, 1])
            cur_type = item.get("type", "short_question")
            type_index = {"mcq": 0, "fill_in_the_blanks": 1, "short_question": 2, "long_question": 3}.get(cur_type, 2)
            item["type"] = cols_meta[0].selectbox(
                "Type",
                ("mcq", "fill_in_the_blanks", "short_question", "long_question"),
                index=type_index,
                key=f"type_{key_prefix}"
            )
            item["points"] = cols_meta[1].number_input(
                "Max Points", min_value=1, value=max(1, item.get("points", 10)), key=f"pts_{key_prefix}"
            )

            # If parent (has parts): show subtotal and render children inside same expander
            if item.get("parts"):
                # compute subtotal from any existing results_map
                subtotal, submax = 0, 0
                def _collect_leaves(parts, prefix):
                    for j, p in enumerate(parts):
                        curp = prefix + [j]
                        if p.get("parts"):
                            yield from _collect_leaves(p["parts"], curp)
                        else:
                            yield (curp, p)
                leaves = list(_collect_leaves(item["parts"], prefix=path))
                for pth, leaf_q in leaves:
                    key = "q-" + "-".join(map(str, pth))
                    r = st.session_state.get("results_map", {}).get(key, {})
                    subtotal += int(r.get("score", 0))
                    submax += int(leaf_q.get("points", 0))

                st.markdown(f"**Subtotal:** {subtotal} / {submax}")
                st.write("")  # spacer
                # Render children inside same expander (nesting visually)
                display_questions(item["parts"], path_prefix=path, depth=depth+1)

            else:
                # Leaf: answer + rubric UI
                item["answer"] = st.text_area("Student Answer", value=item.get("answer", ""), key=f"ans_{key_prefix}", height=100)
                st.markdown("**Rubric (Optional)**")
                r1, r2 = st.columns(2)
                item["rubric_cols"] = r1.number_input("Columns", 2, 5, 3, key=f"rub_cols_{key_prefix}")
                item["suggested_criteria"] = r2.text_input(
                    "Suggest Criteria (comma-separated or 'auto')",
                    value=item.get("suggested_criteria", ""),
                    key=f"sug_crit_{key_prefix}"
                )

                btn_cols = st.columns([1, 1, 1])
                btn_cols[0].button("🧠 Generate AI Rubric", key=f"rub_{key_prefix}", on_click=generate_rubric_for_question, args=(path,))
                btn_cols[1].button("🔁 Regenerate Rubric", key=f"regrub_{key_prefix}", on_click=regenerate_rubric_for_question, args=(path,))
                if item.get("rubric"):
                    btn_cols[2].button("🗑️ Delete Rubric", key=f"delrub_{key_prefix}", on_click=delete_rubric_for_question, args=(path,))

                if item.get("rubric"):
                    st.markdown("**Editable Rubric Table:**")
                    item["rubric"] = st.data_editor(
                        item["rubric"],
                        num_rows="dynamic",
                        use_container_width=True,
                        key=f"edit_rub_{key_prefix}"
                    )




# ================== ADD / CLEAR BUTTONS ==================
st.markdown("---")
cols = st.columns([1, 1, 3])
cols[0].button("➕ Add Top-Level Question", on_click=lambda: st.session_state.exam_questions.append({
    "question": "",
    "answer": "",
    "points": 10,
    "type": "short_question",
    "parts": [],
    "rubric": None,
    "rubric_cols": 3,
    "suggested_criteria": ""
}))
cols[1].button("🧹 Clear All", on_click=lambda: st.session_state.pop("exam_questions", None))

# Bulk-rubric controls
bulk_cols = st.columns([2,1,1,1])
bulk_cols[0].button("🧩 Generate rubrics for ALL leaf questions", on_click=lambda: generate_rubrics_for_all(force=False, columns=3))
bulk_cols[1].checkbox("Force overwrite existing rubrics", key="bulk_force", value=False)
bulk_cols[2].number_input("Columns (for bulk)", min_value=2, max_value=5, value=3, key="bulk_cols")
bulk_cols[3].number_input("Max workers", min_value=1, max_value=12, value=6, key="bulk_workers")

display_questions(st.session_state.exam_questions)

# ================== GRADE BUTTON ==================
st.markdown("---")
if st.button("🚀 Grade Exam", type="primary", use_container_width=True):
    if not st.session_state.exam_questions:
        st.warning("Please add at least one question before grading.")
    else:
        with st.spinner("Grading in progress..."):
            try:
                payload = {
                    "strictness_level": strictness_instruction,
                    "course_material": extract_text_from_files(uploaded_files),
                    "additional_instructions": additional_instructions,
                    "student_exam": st.session_state.exam_questions
                }
                response = requests.post(BACKEND_URL_GRADE, json=payload, timeout=300)
                response.raise_for_status()
                # backend now returns a mapping of path-keys -> result objects
                st.session_state.results_map = response.json()
                st.success("✅ Grading complete!")
            except requests.exceptions.RequestException as e:
                st.error(f"Failed to connect to backend: {e}")

# ================== RESULTS DISPLAY ==================
import json
from datetime import datetime

def yield_leaves_with_paths(parts, prefix=[]):
    """Yield (path_list, question_dict) for every leaf node."""
    for idx, part in enumerate(parts):
        path = prefix + [idx]
        if part.get('parts'):
            yield from yield_leaves_with_paths(part['parts'], prefix=path)
        else:
            yield (path, part)

def calc_total_from_tree(questions, results_map):
    total, max_total = 0, 0
    for path, q in yield_leaves_with_paths(questions, prefix=[]):
        key = "q-" + "-".join(map(str, path))
        res = results_map.get(key, {})
        total += int(res.get("score", 0))
        max_total += int(q.get("points", 0))
    return total, max_total

def render_results_tree(questions, results_map, prefix=[]):
    """
    Walk the question tree and render results for leaves.
    Containers show a subtotal line and their children inside an expander.
    """
    for idx, q in enumerate(questions):
        cur_path = prefix + [idx]
        if q.get("parts"):
            # compute subtotal for this container
            subtotal, submax = 0, 0
            for p, leaf_q in yield_leaves_with_paths(q["parts"], prefix=cur_path):
                k = "q-" + "-".join(map(str, p))
                r = results_map.get(k, {})
                subtotal += int(r.get("score", 0))
                submax += int(leaf_q.get("points", 0))
            with st.expander(f"{q.get('question','(Parent)')} — Subtotal: {subtotal}/{submax}", expanded=False):
                render_results_tree(q["parts"], results_map, prefix=cur_path)
        else:
            key = "q-" + "-".join(map(str, cur_path))
            res = results_map.get(key)
            if not res:
                # create default entry so UI has something to edit
                res = {
                    "score": 0,
                    "feedback": "",
                    "student_answer": q.get("answer", ""),
                    "max_score": q.get("points", 0),
                    "question": q.get("question", "")
                }
                results_map[key] = res

            # Show leaf result with editable controls
            with st.container():
                st.markdown(f"**Question:** {res.get('question', q.get('question',''))}")
                st.markdown("**Student Answer:**")
                st.info(res.get("student_answer", q.get("answer", "")))
                st.markdown("**AI Feedback (editable):**")
                fb_key = f"fb_{key}"
                score_key = f"score_{key}"
                new_fb = st.text_area("Feedback", value=res.get("feedback", ""), key=fb_key, height=120)
                new_score = st.number_input(
                    "Assigned Score",
                    min_value=0,
                    max_value=int(res.get("max_score", q.get("points", 0))),
                    value=int(res.get("score", 0)),
                    step=1,
                    key=score_key
                )
                # update results_map if changed
                if new_fb != res.get("feedback") or int(new_score) != int(res.get("score", 0)):
                    results_map[key]["feedback"] = new_fb
                    results_map[key]["score"] = int(new_score)

# show results if we have them
if "results_map" in st.session_state:
    # defensive: handle error-style responses
    rm = st.session_state.get("results_map")
    if isinstance(rm, dict) and "error" in rm and len(rm)==1:
        st.error(f"Backend returned error: {rm.get('error')}")
    else:
        st.header("📊 Grading Results (editable)")
        results_map = st.session_state.results_map

        # compute totals and render
        total_before, max_before = calc_total_from_tree(st.session_state.exam_questions, results_map)
        render_results_tree(st.session_state.exam_questions, results_map)
        total_after, max_after = calc_total_from_tree(st.session_state.exam_questions, results_map)
        st.subheader(f"🏆 Final Score (after edits): {total_after} / {max_after}")

        # Export and Clear
        c1, c2 = st.columns([1,1])
        export_payload = {
            "exported_at": datetime.now().isoformat(),
            "final_score": {"score": total_after, "max_score": max_after},
            "results_map": results_map,
            "exam_questions_snapshot": st.session_state.exam_questions
        }
        json_str = json.dumps(export_payload, indent=2, ensure_ascii=False)

        c1.download_button("⬇️ Download edited results (JSON)", json_str, file_name="edited_results.json", mime="application/json")
        if c2.button("🗑️ Clear Results"):
            st.session_state.pop("results_map", None)
            st.experimental_rerun()
